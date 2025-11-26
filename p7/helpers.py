"""Helper functions for internal API calls and validation."""

import os
import json
import mimetypes
from pathlib import Path
from typing import Optional
from io import BytesIO
import requests
from django.http import JsonResponse
from pypdf import PdfReader
from docx import Document
from pptx import Presentation
from openpyxl import load_workbook


def validate_internal_auth(x_internal_auth: str) -> JsonResponse | None:
    """
    Validate the internal auth header value. Returns JsonResponse on failure, otherwise None.

    params:
        x_internal_auth (str): The value of the x-internal-auth header to validate.
    """
    if x_internal_auth != os.getenv("INTERNAL_API_KEY"):
        return JsonResponse(
            {"error": "Unauthorized - invalid x-internal-auth"}, status=401
        )
    return None


def fetch_api(url, headers, data):
    """
    Fetches a url with provided headers and data. Contains generic error handling

    params:
        url (str): The URL to fetch.
        headers (dict): The headers to include in the request.
        data (dict): The JSON body to include in the request.
    """
    response = requests.post(url, headers=headers, json=data, timeout=10)
    if not response.ok:
        return JsonResponse(
            {"error": "Failed to fetch files", "details": response.json()},
            status=response.status_code,
        )
    return response

def smart_extension(provider: str, name: str, mime: Optional[str] = None) -> str:
    """
    Determine the file extension based on provider, filename, and MIME type.

    Rules:
    - Dotfiles like ".gitignore" => no extension
    - Names with multiple leading dots (e.g. "..docx", "...pdf") should still yield a real extension
    - Preserve compression combos like ".tar.gz"
    - Fallback to MIME (and Google Drive pseudo-types) when needed
    """
    # known compression endings
    compressed_file_extensions = {'.gz', '.bz2', '.xz', '.zst', '.lz', '.lzma', '.br'}

    # some common modern extensions aren't in Python's built-in mimetypes table
    extra_known_extensions = {
        '.docx', '.xlsx', '.pptx', '.webp', '.md', '.json', '.yaml', '.yml', '.toml',
        '.7z', '.rar', '.heic', '.heif', '.svg', '.ts', '.tsx', '.jsx', '.ipynb',
        '.csv', '.tsv', '.parquet', '.rtf', '.odt', '.ods', '.odp', '.epub',
        '.ttf', '.otf', '.woff', '.woff2'
    }

    known_file_extensions = set(mimetypes.types_map) \
                            | compressed_file_extensions \
                            | extra_known_extensions \
                            | downloadable_file_extensions()

    google_file_extensions = {}
    if provider == "google":
        google_file_extensions = {
            "application/vnd.google-apps.document": ".gdoc",
            "application/vnd.google-apps.spreadsheet": ".gsheet",
            "application/vnd.google-apps.presentation": ".gslides",
            "application/vnd.google-apps.drawing": ".gdraw",
            "application/vnd.google-apps.form": ".gform",
            "application/vnd.google-apps.fusiontable": ".gtable",
            "application/vnd.google-apps.map": ".gmap",
            "application/vnd.google-apps.script": ".gscript",
            "application/vnd.google-apps.site": ".gsite",
            "application/vnd.google-apps.jam": ".gjam",
        }

    filename = (name or "").strip()
    core = filename.lstrip(".")  # ignore any number of leading dots for extension detection

    # If nothing remains after stripping, or there's no further dot, it's a dotfile/plain name.
    if not core:
        pass  # fall through to MIME fallback below
    elif "." not in core:
        # Handle names like "..docx" / "...pdf":
        # if the remaining token itself is a known ext, use it.
        candidate = f".{core.lower()}"
        if candidate in known_file_extensions:
            return candidate
        # otherwise fall through to MIME fallback
    else:
        # Normal path: split and examine trailing parts
        parts = core.lower().split(".")  # ['file','tar','gz'] or ['file','txt'] or ['','docx']
        last = f".{parts[-1]}"

        # Preserve combos like ".tar.gz" or ".csv.gz"
        if last in compressed_file_extensions and len(parts) >= 2:
            penult = f".{parts[-2]}"
            if penult in known_file_extensions:
                return f"{penult}{last}"

        # Single extension case
        if last in known_file_extensions:
            return last

    # Fallback to MIME type
    if mime:
        ext = mimetypes.guess_extension(mime)
        if ext:
            return ext.lower()
        # Google Drive pseudo-MIME types
        if provider == "google":
            return google_file_extensions.get(mime, "")

    return ""


def downloadable_file_extensions() -> set[str]:
    """Return a set of file extensions considered downloadable."""

    file_extensions_path = Path(__file__).resolve().parent / "json/downloadable_txt_extensions.json"

    with file_extensions_path.open("r", encoding="utf-8") as fh:
        downloadable_text_extensions = set(json.load(fh))

    google_file_extensions = {
        ".gdoc",
        ".gsheet",
        ".gslides",
        ".gdraw",
        ".gform",
        ".gtable",
        ".gmap",
        ".gscript",
        ".gsite",
        ".gjam",
    }

    other_file_extensions = {
        ".pdf",
    #   ".doc", # legacy Word format (unsupported)
    #   ".ppt",  # legacy PowerPoint format (unsupported)
    #   ".xls",  # legacy Excel format (unsupported)
        ".docx",
        ".pptx",
        ".xlsx",
    }

    return downloadable_text_extensions | google_file_extensions | other_file_extensions

def parse_file_content(content_bytes: bytes, file) -> str | None:
    """
    Parse file content to extract text from different file types.

    params:
        content (bytes): The raw file content in bytes.
    """

    content_bytes_decoded = content_bytes.decode("utf-8-sig", errors="ignore").strip()
    content_bytes = BytesIO(content_bytes)

    if content_bytes_decoded:
        match file.extension:
            case ".pdf":
                try:
                    reader = PdfReader(content_bytes)
                    return "\n".join(page.extract_text() or "" for page in reader.pages)
                except RuntimeError as e:
                    print(f"Failed to parse pdf: {e}")
                    return None
            case ".docx":
                try:
                    doc = Document(content_bytes)
                    return "\n".join(p.text for p in doc.paragraphs)
                except RuntimeError as e:
                    print(f"Failed to parse docx: {e}")
                    return None
            case ".pptx":
                try:
                    prs = Presentation(content_bytes)
                    slide_text = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if not shape.has_text_frame:
                                continue
                            for paragraph in shape.text_frame.paragraphs:
                                for run in paragraph.runs:
                                    slide_text.append(run.text)
                    return "\n".join(slide_text)
                except RuntimeError as e:
                    print(f"Failed to parse pptx: {e}")
                    return None
            case ".xlsx" | ".gsheet":
                try:
                    wb = load_workbook(content_bytes, data_only=True)
                    all_sheets_text = []
                    for ws in wb.worksheets:
                        rows = [list(row) for row in ws.iter_rows(values_only=True)]
                        sheet_text = "\n".join(
                            "\t".join(str(cell) if cell is not None else "" for cell in row)
                            for row in rows
                        )
                        if sheet_text.strip():  # Only include non-empty sheets
                            all_sheets_text.append(f"\n{sheet_text}")
                    return "\n\n".join(all_sheets_text)
                except RuntimeError as e:
                    print(f"Failed to parse xlsx: {e}")
                    return None
            case _: # Default: try to decode as UTF-8 text
                try:
                    return content_bytes_decoded
                except RuntimeError as e:
                    print(f"Failed to decode: {e}")
                    return None

    return None
